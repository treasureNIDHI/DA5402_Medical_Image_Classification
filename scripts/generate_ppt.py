"""
Generate the project presentation (PPTX) for evaluation.
Run: venv/bin/python scripts/generate_ppt.py
Output: docs/DA5402_Medical_Imaging_MLOps.pptx
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt, Emu

OUT_FILE = Path("docs/DA5402_Medical_Imaging_MLOps.pptx")

# ── Brand colours ──────────────────────────────────────────────────────────
BLUE       = RGBColor(0x36, 0x54, 0xC8)   # primary
DARK_BLUE  = RGBColor(0x1E, 0x2A, 0x6E)   # header bg
PURPLE     = RGBColor(0x76, 0x4A, 0xBC)   # accent
GREEN      = RGBColor(0x10, 0xB9, 0x81)   # success
ORANGE     = RGBColor(0xF5, 0x9E, 0x0B)   # warning
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF3, 0xF4, 0xF6)
DARK_GRAY  = RGBColor(0x37, 0x41, 0x51)
RED        = RGBColor(0xEF, 0x44, 0x44)

W, H = Inches(13.33), Inches(7.5)          # 16:9 widescreen


# ── Helpers ────────────────────────────────────────────────────────────────

def new_slide(prs: Presentation, layout_idx: int = 6):
    """Add a blank slide."""
    layout = prs.slide_layouts[layout_idx]
    return prs.slides.add_slide(layout)


def rect(slide, x, y, w, h, fill: RGBColor | None = None, alpha: int | None = None):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.line.fill.background()
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    return shape


def txbox(slide, text: str, x, y, w, h,
          size=18, bold=False, color=DARK_GRAY, align=PP_ALIGN.LEFT,
          wrap=True, italic=False):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    box.word_wrap = wrap
    tf = box.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return box


def header_bar(slide, title: str, subtitle: str = ""):
    rect(slide, 0, 0, 13.33, 1.3, fill=DARK_BLUE)
    txbox(slide, title, 0.3, 0.1, 12, 0.7, size=28, bold=True, color=WHITE)
    if subtitle:
        txbox(slide, subtitle, 0.3, 0.78, 12, 0.45, size=15, color=RGBColor(0xA5, 0xB4, 0xFC))


def bullet_box(slide, items: list[str], x, y, w, h,
               size=16, color=DARK_GRAY, bullet="•"):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    box.word_wrap = True
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = f"{bullet}  {item}"
        run.font.size = Pt(size)
        run.font.color.rgb = color
        p.space_after = Pt(4)


def metric_card(slide, label: str, value: str, x, y,
                bg=BLUE, val_color=WHITE, lbl_color=RGBColor(0xB4, 0xC6, 0xFF)):
    rect(slide, x, y, 2.8, 1.3, fill=bg)
    txbox(slide, value, x+0.1, y+0.08, 2.6, 0.7, size=32, bold=True,
          color=val_color, align=PP_ALIGN.CENTER)
    txbox(slide, label, x+0.1, y+0.78, 2.6, 0.42, size=13,
          color=lbl_color, align=PP_ALIGN.CENTER)


def divider(slide, y, color=BLUE):
    line = slide.shapes.add_shape(1, Inches(0.3), Inches(y), Inches(12.73), Inches(0.03))
    line.fill.solid()
    line.fill.fore_color.rgb = color
    line.line.fill.background()


# ══════════════════════════════════════════════════════════════════════════
# SLIDES
# ══════════════════════════════════════════════════════════════════════════

def slide_title(prs):
    sl = new_slide(prs)
    # Full background gradient simulation
    rect(sl, 0, 0, 13.33, 7.5, fill=DARK_BLUE)
    rect(sl, 0, 5.5, 13.33, 2.0, fill=RGBColor(0x12, 0x1A, 0x50))
    # Accent bar
    rect(sl, 0, 4.7, 13.33, 0.08, fill=PURPLE)

    txbox(sl, "Medical Image Classification", 1, 1.0, 11.33, 1.1,
          size=38, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txbox(sl, "End-to-End MLOps System", 1, 2.0, 11.33, 0.8,
          size=28, color=RGBColor(0xA5, 0xB4, 0xFC), align=PP_ALIGN.CENTER)
    divider(sl, 3.05, color=PURPLE)
    txbox(sl, "Pneumonia Detection  ·  Brain Tumor Classification", 1, 3.2, 11.33, 0.6,
          size=18, color=RGBColor(0xC7, 0xD2, 0xFE), align=PP_ALIGN.CENTER)
    txbox(sl, "DA5402 · April 2026", 1, 5.7, 11.33, 0.5,
          size=14, color=RGBColor(0x6B, 0x7A, 0xAA), align=PP_ALIGN.CENTER)

    # Tool badges
    tools = ["DVC", "MLflow", "FastAPI", "Docker", "Prometheus", "Grafana", "Airflow"]
    for i, t in enumerate(tools):
        bx = 0.5 + i * 1.82
        rect(sl, bx, 4.15, 1.6, 0.42, fill=RGBColor(0x2D, 0x3A, 0x8C))
        txbox(sl, t, bx, 4.18, 1.6, 0.38, size=12, bold=True,
              color=RGBColor(0xA5, 0xB4, 0xFC), align=PP_ALIGN.CENTER)


def slide_problem(prs):
    sl = new_slide(prs)
    header_bar(sl, "Problem Statement", "Two clinical AI challenges requiring production-grade solutions")
    rect(sl, 0, 1.3, 13.33, 6.2, fill=LIGHT_GRAY)

    # Left card
    rect(sl, 0.4, 1.6, 5.9, 5.5, fill=WHITE)
    rect(sl, 0.4, 1.6, 5.9, 0.5, fill=BLUE)
    txbox(sl, "🫁  Pneumonia Detection", 0.5, 1.63, 5.7, 0.44,
          size=15, bold=True, color=WHITE)
    bullet_box(sl, [
        "Chest X-ray binary classification",
        "Classes: NORMAL vs PNEUMONIA",
        "5,856 images (Kaggle, CC BY 4.0)",
        "Pediatric patients, Guangzhou Medical Center",
        "Train/Val/Test: 5216 / 16 / 624",
        "Class imbalance: 1:2.9 (NORMAL:PNEUMONIA)",
        "Goal: Automate triage & reduce radiologist load",
    ], 0.55, 2.25, 5.6, 4.6, size=14)

    # Right card
    rect(sl, 7.0, 1.6, 5.9, 5.5, fill=WHITE)
    rect(sl, 7.0, 1.6, 5.9, 0.5, fill=PURPLE)
    txbox(sl, "🧠  Brain Tumor Classification", 7.1, 1.63, 5.7, 0.44,
          size=15, bold=True, color=WHITE)
    bullet_box(sl, [
        "Brain MRI 4-class classification",
        "Classes: glioma / meningioma / notumor / pituitary",
        "7,023 images (Kaggle)",
        "Training: 5,712  |  Testing: 1,311",
        "Near-balanced class distribution",
        "Goal: Early detection to support treatment decisions",
        "Achieves 92.7% test accuracy",
    ], 7.15, 2.25, 5.6, 4.6, size=14)

    txbox(sl, "⚠  Not for clinical use — research & educational purposes only",
          1.5, 6.95, 10.33, 0.4, size=11, italic=True,
          color=RGBColor(0x9C, 0xA3, 0xAF), align=PP_ALIGN.CENTER)


def slide_architecture(prs):
    sl = new_slide(prs)
    header_bar(sl, "System Architecture", "7-layer design with strict separation of concerns")
    rect(sl, 0, 1.3, 13.33, 6.2, fill=LIGHT_GRAY)

    layers = [
        ("User Layer",        "Browser SPA (Nginx :80) — 4-tab dashboard",               BLUE,      0.4),
        ("Inference Layer",   "FastAPI :8001 · predict.py · health.py · Prometheus MW",  PURPLE,    0.4),
        ("Model Layer",       "pneumonia_resnet50.pt + brain_resnet50.pt  (baked in Docker image)", RGBColor(0x0E, 0x78, 0x90), 0.4),
        ("MLOps Layer",       "MLflow (tracking + registry) · DVC (pipeline + versioning) · Git",  RGBColor(0x05, 0x7A, 0x55), 0.4),
        ("Monitoring Layer",  "Prometheus :9090 · Grafana :3001 · monitor.py (drift detection)",    ORANGE,    0.4),
        ("Training Layer",    "train.py · train_brain.py · optimization.py  [offline, not in inference path]", RGBColor(0x92, 0x40, 0x0E), 0.4),
        ("Orchestration",     "Airflow :8080 (DAG scheduler) · GitHub Actions (CI/CD on push)",      DARK_GRAY, 0.4),
    ]
    arrows = ["↓ REST /predict", "↓ loads .pt at startup", "↓ logs runs", "↓ scrapes /metrics", "↓ triggered offline", ""]

    for i, (name, desc, color, _) in enumerate(layers):
        y = 1.45 + i * 0.77
        rect(sl, 0.35, y, 2.5, 0.63, fill=color)
        txbox(sl, name, 0.38, y+0.05, 2.44, 0.55, size=12, bold=True,
              color=WHITE, align=PP_ALIGN.CENTER)
        rect(sl, 2.85, y, 9.85, 0.63, fill=WHITE)
        txbox(sl, desc, 2.95, y+0.08, 9.65, 0.52, size=12, color=DARK_GRAY)
        if i < len(arrows) - 1:
            txbox(sl, arrows[i], 0.55, y+0.63, 2.1, 0.22,
                  size=9, color=RGBColor(0x9C, 0xA3, 0xAF), align=PP_ALIGN.CENTER)

    txbox(sl, "Key: Models baked into Docker image — inference requires zero external dependencies",
          0.35, 6.98, 12.63, 0.38, size=11, bold=True,
          color=BLUE, align=PP_ALIGN.CENTER)


def slide_data_pipeline(prs):
    sl = new_slide(prs)
    header_bar(sl, "Data Engineering Pipeline", "13-stage DVC pipeline · fully reproducible · MD5-checksummed")
    rect(sl, 0, 1.3, 13.33, 6.2, fill=LIGHT_GRAY)

    stages = [
        ("1 Ingestion",    BLUE),   ("2 Validate",    GREEN),  ("3 Preprocess",  BLUE),
        ("4 Feature Eng",  PURPLE), ("5 EDA",         PURPLE), ("6 Split",       BLUE),
        ("7 Train Pneumo", ORANGE), ("8 Train Brain",  ORANGE), ("9 Evaluate",    GREEN),
        ("10 Optimize",    PURPLE), ("11 Experiments", BLUE),  ("12 Monitor",    GREEN),
        ("13 Health Chk",  GREEN),
    ]
    cols, rows = 4, 4
    for i, (name, color) in enumerate(stages):
        col, row = i % cols, i // cols
        bx = 0.45 + col * 3.1
        by = 1.55 + row * 1.25
        rect(sl, bx, by, 2.7, 0.8, fill=color)
        txbox(sl, name, bx, by, 2.7, 0.8, size=13, bold=True,
              color=WHITE, align=PP_ALIGN.CENTER)
        if col < cols - 1 and i < len(stages) - 1:
            txbox(sl, "→", bx+2.72, by+0.22, 0.35, 0.35,
                  size=18, bold=True, color=DARK_GRAY, align=PP_ALIGN.CENTER)

    txbox(sl, "Airflow DAG chains all stages · dvc dag shows dependency graph · dvc repro runs end-to-end",
          0.4, 6.6, 12.53, 0.4, size=12, color=DARK_GRAY, align=PP_ALIGN.CENTER)
    rect(sl, 0.35, 6.55, 12.63, 0.5, fill=WHITE)
    txbox(sl, "✓  All 13 stages completed  ·  13,056 images processed  ·  ~45 min total runtime",
          0.4, 6.58, 12.53, 0.4, size=13, bold=True, color=GREEN, align=PP_ALIGN.CENTER)


def slide_model(prs):
    sl = new_slide(prs)
    header_bar(sl, "Model Development", "ResNet-50 (ImageNet pretrained) — fine-tuned for medical imaging")
    rect(sl, 0, 1.3, 13.33, 6.2, fill=LIGHT_GRAY)

    # Architecture justification table
    rect(sl, 0.35, 1.45, 7.0, 0.45, fill=DARK_BLUE)
    for label, x in [("Model", 0.4), ("Test Acc", 2.5), ("Latency", 3.9), ("Size", 5.1), ("Decision", 5.9)]:
        txbox(sl, label, x, 1.47, 1.4, 0.4, size=12, bold=True, color=WHITE)

    rows_data = [
        ("ResNet-50 ✓", "81.4% / 92.7%", "~5ms", "90MB", "✅ SELECTED", GREEN),
        ("ResNet-34",   "~89%",           "~4ms", "80MB", "Lower acc",   DARK_GRAY),
        ("MobileNet V2","~87%",           "~2ms", "14MB", "5% gap",      DARK_GRAY),
        ("EfficientNet B4","~94%",        ">200ms","75MB","Too slow",    ORANGE),
        ("ViT",         "~93%",           "~15ms","330MB","Too large",   ORANGE),
    ]
    for i, (name, acc, lat, sz, dec, dc) in enumerate(rows_data):
        by = 1.92 + i * 0.52
        bg = WHITE if i % 2 == 0 else LIGHT_GRAY
        rect(sl, 0.35, by, 7.0, 0.5, fill=bg)
        bold = (i == 0)
        for val, x in [(name, 0.4), (acc, 2.5), (lat, 3.9), (sz, 5.1), (dec, 5.9)]:
            txbox(sl, val, x, by+0.08, 1.45, 0.38, size=12, bold=bold,
                  color=dc if i == 0 else DARK_GRAY)

    # Right: Training config
    rect(sl, 7.65, 1.45, 5.3, 0.45, fill=DARK_BLUE)
    txbox(sl, "Training Configuration", 7.7, 1.47, 5.2, 0.4, size=14, bold=True, color=WHITE)
    configs = [
        ("Optimizer", "Adam"),
        ("Learning Rate", "1e-4"),
        ("Epochs", "5"),
        ("Batch Size", "32"),
        ("Loss", "CrossEntropyLoss"),
        ("Device", "MPS (Apple M-series)"),
        ("Pretrained", "ImageNet weights"),
    ]
    for i, (k, v) in enumerate(configs):
        by = 1.92 + i * 0.52
        bg = WHITE if i % 2 == 0 else LIGHT_GRAY
        rect(sl, 7.65, by, 5.3, 0.5, fill=bg)
        txbox(sl, k, 7.7, by+0.08, 2.5, 0.38, size=12, bold=True, color=DARK_GRAY)
        txbox(sl, v, 10.2, by+0.08, 2.7, 0.38, size=12, color=BLUE)


def slide_results(prs):
    sl = new_slide(prs)
    header_bar(sl, "Model Results", "Trained on MPS (Apple Silicon) · 5 epochs · ResNet-50")
    rect(sl, 0, 1.3, 13.33, 6.2, fill=LIGHT_GRAY)

    # Metric cards
    cards = [
        ("Pneumonia\nTest Accuracy", "81.4%",  BLUE,   0.4),
        ("Brain Tumor\nTest Accuracy","92.7%", PURPLE, 3.5),
        ("Brain Tumor\nF1 Score",    "0.925",  GREEN,  6.6),
        ("Inference\nLatency (CPU)", "~5ms",   ORANGE, 9.7),
    ]
    for label, value, color, x in cards:
        metric_card(sl, label, value, x, 1.6, bg=color)

    # Epoch tables side by side
    rect(sl, 0.35, 3.2, 5.9, 0.45, fill=DARK_BLUE)
    txbox(sl, "Pneumonia — Validation per Epoch", 0.4, 3.22, 5.8, 0.4,
          size=13, bold=True, color=WHITE)
    pneu = [("0","0.9375","0.9373"),("1","0.5625","0.4589"),("2","0.6250","0.5636"),
            ("3","1.0000","1.0000"),("4","1.0000","1.0000")]
    for i, (ep, acc, f1) in enumerate(pneu):
        by = 3.67 + i * 0.49
        bg = WHITE if i % 2 == 0 else LIGHT_GRAY
        rect(sl, 0.35, by, 5.9, 0.48, fill=bg)
        for val, x in [(f"Epoch {ep}", 0.4), (f"Acc {acc}", 2.3), (f"F1 {f1}", 4.2)]:
            txbox(sl, val, x, by+0.08, 2.0, 0.36, size=12,
                  color=GREEN if acc == "1.0000" else DARK_GRAY)

    rect(sl, 6.85, 3.2, 5.9, 0.45, fill=DARK_BLUE)
    txbox(sl, "Brain Tumor — Validation per Epoch", 6.9, 3.22, 5.8, 0.4,
          size=13, bold=True, color=WHITE)
    brain = [("0","0.9268","0.9267"),("1","0.9625","0.9626"),("2","0.9643","0.9643"),
             ("3","0.9688","0.9689"),("4","0.9688","0.9686")]
    for i, (ep, acc, f1) in enumerate(brain):
        by = 3.67 + i * 0.49
        bg = WHITE if i % 2 == 0 else LIGHT_GRAY
        rect(sl, 6.85, by, 5.9, 0.48, fill=bg)
        for val, x in [(f"Epoch {ep}", 6.9), (f"Acc {acc}", 8.8), (f"F1 {f1}", 10.7)]:
            txbox(sl, val, x, by+0.08, 2.0, 0.36, size=12,
                  color=PURPLE if float(acc) >= 0.96 else DARK_GRAY)


def slide_mlops(prs):
    sl = new_slide(prs)
    header_bar(sl, "MLOps Implementation", "Full lifecycle: versioning · tracking · monitoring · retraining")
    rect(sl, 0, 1.3, 13.33, 6.2, fill=LIGHT_GRAY)

    tools = [
        ("DVC", "Pipeline + Data Versioning",
         ["13-stage pipeline in dvc.yaml", "MD5 checksums in dvc.lock",
          "dvc dag · dvc repro · dvc status",  "Data, splits & models versioned"], BLUE),
        ("MLflow", "Experiment Tracking",
         ["3 experiments tracked", "Params: LR, optimizer, epochs",
          "Metrics: acc, F1, loss per epoch", "Artifacts: reports, confusion matrix",
          "Model registry: 2 models in Staging"], PURPLE),
        ("Prometheus\n+ Grafana", "Observability",
         ["Metrics: requests, latency, errors",  "Histogram: p50/p95/p99 latency",
          "Auto-provisioned Grafana dashboard", "Alert rules in alert_rules.yml"], ORANGE),
        ("Airflow", "Orchestration",
         ["DAG: medical_imaging_pipeline", "All 13 DVC stages wired",
          "Retry logic + web UI", "Accessible at :8080"], GREEN),
    ]

    for i, (name, subtitle, bullets, color) in enumerate(tools):
        col = i % 2
        row = i // 2
        bx = 0.4 + col * 6.5
        by = 1.5 + row * 2.75
        rect(sl, bx, by, 6.0, 0.55, fill=color)
        txbox(sl, f"{name}  —  {subtitle}", bx+0.1, by+0.07, 5.8, 0.45,
              size=14, bold=True, color=WHITE)
        rect(sl, bx, by+0.55, 6.0, 2.1, fill=WHITE)
        bullet_box(sl, bullets, bx+0.15, by+0.6, 5.7, 2.0, size=13)


def slide_deployment(prs):
    sl = new_slide(prs)
    header_bar(sl, "Containerization & Deployment", "Self-contained Docker image · models baked in · no retraining needed")
    rect(sl, 0, 1.3, 13.33, 6.2, fill=LIGHT_GRAY)

    # Docker image facts
    rect(sl, 0.35, 1.5, 5.8, 0.5, fill=DARK_BLUE)
    txbox(sl, "Docker Image — medical-imaging:latest", 0.4, 1.52, 5.7, 0.45,
          size=14, bold=True, color=WHITE)
    img_facts = [
        ("Base image", "python:3.11-slim"),
        ("Total size",  "2.41 GB (PyTorch CPU + 2×90MB models)"),
        ("Models",      "✓ Baked in at build time (COPY *.pt)"),
        ("Runtime user","Non-root appuser (uid 1000)"),
        ("Healthcheck", "curl /healthz every 30s"),
        ("Port",        "8001"),
    ]
    for i, (k, v) in enumerate(img_facts):
        by = 2.02 + i * 0.52
        bg = WHITE if i % 2 == 0 else LIGHT_GRAY
        rect(sl, 0.35, by, 5.8, 0.5, fill=bg)
        txbox(sl, k, 0.4, by+0.08, 2.0, 0.38, size=12, bold=True, color=DARK_GRAY)
        txbox(sl, v, 2.42, by+0.08, 3.7, 0.38, size=12,
              color=GREEN if "✓" in v else BLUE)

    # Compose profiles
    rect(sl, 6.6, 1.5, 6.35, 0.5, fill=DARK_BLUE)
    txbox(sl, "Docker Compose Profiles", 6.65, 1.52, 6.25, 0.45,
          size=14, bold=True, color=WHITE)
    profiles = [
        ("--profile inference", "API · Frontend · Prometheus · Grafana", GREEN),
        ("--profile training",  "MLflow · PostgreSQL · Airflow",         ORANGE),
        ("(no profile)",        "Everything — full stack",               PURPLE),
    ]
    for i, (cmd, services, color) in enumerate(profiles):
        by = 2.02 + i * 0.9
        rect(sl, 6.6, by, 6.35, 0.82, fill=WHITE)
        rect(sl, 6.6, by, 0.12, 0.82, fill=color)
        txbox(sl, cmd, 6.8, by+0.04, 5.95, 0.38, size=13, bold=True, color=DARK_BLUE)
        txbox(sl, services, 6.8, by+0.4, 5.95, 0.36, size=12, color=DARK_GRAY)

    # Services table
    rect(sl, 0.35, 5.3, 12.6, 0.45, fill=DARK_BLUE)
    txbox(sl, "Service Endpoints", 0.4, 5.32, 12.5, 0.4, size=13, bold=True, color=WHITE)
    services = [
        ("Frontend",   "http://localhost",       "Web UI — 4 tabs"),
        ("API",        "http://localhost:8001",   "Inference REST API"),
        ("API Docs",   "http://localhost:8001/docs","Swagger UI"),
        ("MLflow",     "http://localhost:5000",   "Experiment tracking"),
        ("Prometheus", "http://localhost:9090",   "Metrics store"),
        ("Grafana",    "http://localhost:3001",   "Dashboards (admin/admin)"),
        ("Airflow",    "http://localhost:8080",   "Pipeline orchestration"),
    ]
    for i, (name, url, desc) in enumerate(services):
        bx = 0.35 + (i % 4) * 3.15
        by = 5.77 + (i // 4) * 0.55
        bg = WHITE if i % 2 == 0 else LIGHT_GRAY
        rect(sl, bx, by, 3.1, 0.5, fill=bg)
        txbox(sl, f"{name}: {url}", bx+0.1, by+0.06, 2.95, 0.2, size=10, bold=True, color=BLUE)
        txbox(sl, desc, bx+0.1, by+0.26, 2.95, 0.2, size=10, color=DARK_GRAY)


def slide_frontend(prs):
    sl = new_slide(prs)
    header_bar(sl, "Frontend — Web Dashboard", "Single-page app · 4 tabs · live API connection · Nginx :80")
    rect(sl, 0, 1.3, 13.33, 6.2, fill=LIGHT_GRAY)

    tabs = [
        ("🔮 Predict", BLUE, [
            "Model selector: Pneumonia vs Brain Tumor",
            "File upload (JPEG/PNG/BMP/TIFF, max 10MB)",
            "Real-time prediction with confidence %",
            "Modality mismatch detection",
            "Low-confidence warning (< 70%)",
            "Inference time displayed",
        ]),
        ("📊 Pipeline", PURPLE, [
            "13-stage DVC DAG visualization",
            "Stage status: completed (green)",
            "Real model metrics: 81.4% / 92.7%",
            "Training configuration summary",
            "Pipeline duration and run date",
        ]),
        ("📈 Monitoring", GREEN, [
            "Calls GET /health every 30 seconds",
            "Live: API status, models loaded",
            "Live: total requests, avg latency",
            "Live: error rate, CPU/memory/disk",
            "Auto-refresh + manual Refresh button",
            "Alert banner on API unreachable",
        ]),
        ("📚 User Manual", ORANGE, [
            "Step-by-step usage guide",
            "Understanding predictions",
            "FAQ for non-technical users",
            "Technical details & API example",
            "Best practices for image upload",
        ]),
    ]
    for i, (name, color, bullets) in enumerate(tabs):
        bx = 0.35 + i * 3.15
        rect(sl, bx, 1.5, 2.95, 0.55, fill=color)
        txbox(sl, name, bx, 1.52, 2.95, 0.5, size=14, bold=True,
              color=WHITE, align=PP_ALIGN.CENTER)
        rect(sl, bx, 2.05, 2.95, 5.1, fill=WHITE)
        bullet_box(sl, bullets, bx+0.1, 2.12, 2.78, 4.95, size=12)


def slide_testing(prs):
    sl = new_slide(prs)
    header_bar(sl, "Testing", "34 / 34 unit tests passed · pytest 9.0.3 · Python 3.13")
    rect(sl, 0, 1.3, 13.33, 6.2, fill=LIGHT_GRAY)

    # Big result
    rect(sl, 0.35, 1.5, 5.5, 2.0, fill=GREEN)
    txbox(sl, "34 / 34", 0.35, 1.55, 5.5, 1.3,
          size=60, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txbox(sl, "TESTS PASSED  ·  100%  ·  17 seconds", 0.35, 2.8, 5.5, 0.55,
          size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Suite breakdown
    suites = [
        ("test_api.py",       "API Endpoints",      "16 tests", "• All HTTP status codes\n• Field presence\n• Confidence range\n• 422 for invalid input"),
        ("test_data.py",      "Data Pipeline",       "7 tests",  "• Validation on empty dirs\n• Preprocessing resize 224×224\n• Corrupt file handling\n• EDA report generation"),
        ("test_inference.py", "Inference & Models",  "11 tests", "• Model load 2-class/4-class\n• Eval mode verification\n• Modality detection\n• predict() with/without model\n• Health check report"),
    ]
    for i, (fname, title, count, details) in enumerate(suites):
        bx = 6.2 + i * 2.35
        by = 1.5
        rect(sl, bx, by, 2.2, 0.5, fill=BLUE)
        txbox(sl, f"{count}", bx, by+0.05, 2.2, 0.42, size=18, bold=True,
              color=WHITE, align=PP_ALIGN.CENTER)
        rect(sl, bx, by+0.5, 2.2, 0.4, fill=DARK_BLUE)
        txbox(sl, title, bx, by+0.52, 2.2, 0.36, size=11, bold=True,
              color=WHITE, align=PP_ALIGN.CENTER)
        rect(sl, bx, by+0.9, 2.2, 2.6, fill=WHITE)
        txbox(sl, fname, bx+0.08, by+0.94, 2.1, 0.3, size=10, italic=True, color=PURPLE)
        txbox(sl, details, bx+0.08, by+1.2, 2.1, 2.2, size=11, color=DARK_GRAY)

    # Acceptance criteria
    rect(sl, 0.35, 3.7, 12.6, 0.45, fill=DARK_BLUE)
    txbox(sl, "Acceptance Criteria — All Met", 0.4, 3.72, 12.5, 0.4,
          size=13, bold=True, color=WHITE)
    criteria = [
        "100% unit test pass rate",
        "Inference latency < 200ms  (actual: ~70ms)",
        "Model size < 100 MB  (actual: 90 MB)",
        "All 13 DVC stages complete",
        "Invalid input → 422 (not 500)",
        "Docker build + smoke test pass",
    ]
    for i, c in enumerate(criteria):
        bx = 0.35 + (i % 3) * 4.2
        by = 4.22 + (i // 3) * 0.72
        rect(sl, bx, by, 4.1, 0.6, fill=WHITE)
        txbox(sl, f"✓  {c}", bx+0.1, by+0.1, 3.9, 0.45, size=12, color=GREEN)


def slide_monitoring(prs):
    sl = new_slide(prs)
    header_bar(sl, "Monitoring & Observability", "Prometheus · Grafana · Drift Detection · Health Checks")
    rect(sl, 0, 1.3, 13.33, 6.2, fill=LIGHT_GRAY)

    blocks = [
        ("Prometheus Metrics", BLUE, [
            "Scrapes /metrics every 15s",
            "http_requests_total (endpoint, method, status)",
            "http_request_duration_seconds (histogram)",
            "Alert rules: API down, latency > 200ms",
        ]),
        ("Grafana Dashboard", PURPLE, [
            "Auto-provisioned via grafana/provisioning/",
            "No manual setup needed",
            "Panels: predictions, latency p50/p95/p99",
            "Request rate by endpoint (timeseries)",
            "Success rate gauge",
        ]),
        ("Drift Detection", ORANGE, [
            "L1 label distribution distance (threshold 0.15)",
            "Feature mean shift (threshold 0.05)",
            "Feature variance shift (threshold 0.25)",
            "Accuracy drop vs baseline (threshold 3%)",
        ]),
        ("Health Checks", GREEN, [
            "3 levels: /healthz · /readyz · health_check.py",
            "Checks: models loadable + inference OK",
            "System: CPU, memory, disk thresholds",
            "Pipeline artifacts presence check",
            "Auto-retraining on drift / perf drop",
        ]),
    ]
    for i, (title, color, bullets) in enumerate(blocks):
        bx = 0.35 + (i % 2) * 6.5
        by = 1.5 + (i // 2) * 2.75
        rect(sl, bx, by, 6.1, 0.52, fill=color)
        txbox(sl, title, bx+0.1, by+0.07, 5.9, 0.42,
              size=14, bold=True, color=WHITE)
        rect(sl, bx, by+0.52, 6.1, 2.15, fill=WHITE)
        bullet_box(sl, bullets, bx+0.15, by+0.58, 5.85, 2.05, size=13)


def slide_software_eng(prs):
    sl = new_slide(prs)
    header_bar(sl, "Software Engineering", "Design principles · Exception handling · Logging · CI/CD")
    rect(sl, 0, 1.3, 13.33, 6.2, fill=LIGHT_GRAY)

    left_items = [
        ("Architecture Diagram", "docs/ARCHITECTURE.md — Mermaid block diagram, 7 layers"),
        ("HLD Document", "docs/HLD.md — design goals, tech choices, data flow, deployment"),
        ("LLD Document", "docs/LLD.md — 8 endpoints, request/response schemas, error codes"),
        ("Loose Coupling", "Frontend ↔ Backend via REST only — API_URL is the single contract"),
        ("Structured Logging", "Per-request INFO/WARNING/ERROR in FastAPI, full traceback on exception"),
        ("Exception Handling", "ValueError → 422, RuntimeError → 503, generic Exception → 500"),
        ("Non-root Docker", "Container runs as appuser uid 1000 — security best practice"),
    ]
    rect(sl, 0.35, 1.5, 6.1, 0.45, fill=DARK_BLUE)
    txbox(sl, "Design & Implementation", 0.4, 1.52, 6.0, 0.4, size=13, bold=True, color=WHITE)
    for i, (k, v) in enumerate(left_items):
        by = 1.97 + i * 0.66
        bg = WHITE if i % 2 == 0 else LIGHT_GRAY
        rect(sl, 0.35, by, 6.1, 0.63, fill=bg)
        txbox(sl, k, 0.45, by+0.04, 5.9, 0.25, size=12, bold=True, color=DARK_BLUE)
        txbox(sl, v, 0.45, by+0.3, 5.9, 0.3, size=11, color=DARK_GRAY)

    right_items = [
        ("GitHub Actions CI", "Syntax check 20 modules · DVC validate · Docker build on push"),
        ("MLproject", "6 entry points: train, evaluate, monitor, health_check, full_pipeline"),
        ("User Manual", "docs/USER_MANUAL.md + in-app tab — for non-technical users"),
        ("CONTRIBUTING.md", "Branch strategy, commit conventions, PR process, test commands"),
        ("RELEASE_CHECKLIST.md","Pre-release quality gates + rollback procedure"),
        ("DATA_SOURCES_AND_BIAS.md","Provenance, bias analysis, mitigation strategies"),
        ("Reproducibility", "Fixed seed=42, git commit in every MLflow run, DVC checksums"),
    ]
    rect(sl, 6.85, 1.5, 6.1, 0.45, fill=DARK_BLUE)
    txbox(sl, "Documentation & Process", 6.9, 1.52, 6.0, 0.4, size=13, bold=True, color=WHITE)
    for i, (k, v) in enumerate(right_items):
        by = 1.97 + i * 0.66
        bg = WHITE if i % 2 == 0 else LIGHT_GRAY
        rect(sl, 6.85, by, 6.1, 0.63, fill=bg)
        txbox(sl, k, 6.95, by+0.04, 5.9, 0.25, size=12, bold=True, color=DARK_BLUE)
        txbox(sl, v, 6.95, by+0.3, 5.9, 0.3, size=11, color=DARK_GRAY)


def slide_summary(prs):
    sl = new_slide(prs)
    rect(sl, 0, 0, 13.33, 7.5, fill=DARK_BLUE)
    rect(sl, 0, 6.5, 13.33, 1.0, fill=RGBColor(0x12, 0x1A, 0x50))
    rect(sl, 0, 1.25, 13.33, 0.06, fill=PURPLE)

    txbox(sl, "Project Summary", 0.5, 0.15, 12.33, 0.9,
          size=34, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    cards = [
        ("🎯 Problem",       "Pneumonia detection\nBrain tumor 4-class\nclinical screening",   BLUE),
        ("🏗️ Architecture",   "7-layer system\nREST-decoupled\nDocker profiles",               PURPLE),
        ("🤖 Models",         "ResNet-50 pretrained\n81.4% / 92.7% acc\n~5ms inference",        RGBColor(0x0E, 0x78, 0x90)),
        ("📊 MLOps",          "DVC · MLflow · Airflow\nPrometheus · Grafana\n13-stage pipeline", RGBColor(0x05, 0x7A, 0x55)),
        ("🐳 Deploy",         "Models baked in\ndocker compose\n2.41GB image",                  ORANGE),
        ("✅ Testing",        "34 / 34 passed\n100% pass rate\n17 sec runtime",                 GREEN),
    ]
    for i, (title, body, color) in enumerate(cards):
        col = i % 3
        row = i // 3
        bx = 0.5 + col * 4.12
        by = 1.55 + row * 2.42
        rect(sl, bx, by, 3.8, 0.5, fill=color)
        txbox(sl, title, bx, by+0.05, 3.8, 0.42, size=14, bold=True,
              color=WHITE, align=PP_ALIGN.CENTER)
        rect(sl, bx, by+0.5, 3.8, 1.82, fill=RGBColor(0x1E, 0x2D, 0x6E))
        txbox(sl, body, bx, by+0.55, 3.8, 1.72, size=14,
              color=RGBColor(0xC7, 0xD2, 0xFE), align=PP_ALIGN.CENTER)

    txbox(sl, "DA5402  ·  Medical Image Classification MLOps  ·  April 2026",
          0.5, 6.6, 12.33, 0.7, size=13,
          color=RGBColor(0x6B, 0x7A, 0xAA), align=PP_ALIGN.CENTER)


# ── Main ───────────────────────────────────────────────────────────────────

def main():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H

    slide_title(prs)
    slide_problem(prs)
    slide_architecture(prs)
    slide_data_pipeline(prs)
    slide_model(prs)
    slide_results(prs)
    slide_mlops(prs)
    slide_deployment(prs)
    slide_frontend(prs)
    slide_testing(prs)
    slide_monitoring(prs)
    slide_software_eng(prs)
    slide_summary(prs)

    OUT_FILE.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT_FILE)
    print(f"✓ Saved {OUT_FILE}  ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
